import os
import json
import sqlite3
from pathlib import Path
from docx import Document
import PyPDF2
import nbformat
from pptx import Presentation

def extraer_texto_docx(archivo):
    """Extrae texto de archivos .docx"""
    try:
        doc = Document(archivo)
        texto = []
        for parrafo in doc.paragraphs:
            if parrafo.text.strip():
                texto.append(parrafo.text)
        return '\n'.join(texto)
    except Exception as e:
        return f"Error procesando DOCX: {e}"

def extraer_texto_pdf(archivo):
    """Extrae texto de archivos .pdf"""
    try:
        texto = []
        with open(archivo, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for pagina in pdf_reader.pages:
                texto.append(pagina.extract_text())
        return '\n'.join(texto)
    except Exception as e:
        return f"Error procesando PDF: {e}"

def extraer_texto_ipynb(archivo):
    """Extrae texto y código de archivos .ipynb"""
    try:
        with open(archivo, 'r', encoding='utf-8') as file:
            notebook = nbformat.read(file, as_version=4)
        
        contenido = []
        for celda in notebook.cells:
            if celda.cell_type == 'markdown':
                contenido.append(f"# {celda.source}")
            elif celda.cell_type == 'code':
                contenido.append(f"```python\n{celda.source}\n```")
                if celda.outputs:
                    for output in celda.outputs:
                        if hasattr(output, 'text'):
                            contenido.append(f"Output: {output.text}")
        
        return '\n'.join(contenido)
    except Exception as e:
        return f"Error procesando IPYNB: {e}"

def extraer_texto_pptx(archivo):
    """Extrae texto de archivos .pptx"""
    try:
        prs = Presentation(archivo)
        texto = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texto.append(shape.text)
        return '\n'.join(texto)
    except Exception as e:
        return f"Error procesando PPTX: {e}"

def extraer_texto_sqlite(archivo):
    """Extrae estructura y datos de archivos .sqlite"""
    try:
        conn = sqlite3.connect(archivo)
        cursor = conn.cursor()
        
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tablas = cursor.fetchall()
        
        contenido = []
        for tabla in tablas:
            nombre_tabla = tabla[0]
            contenido.append(f"Tabla: {nombre_tabla}")
            
            cursor.execute(f"PRAGMA table_info({nombre_tabla})")
            columnas = cursor.fetchall()
            contenido.append("Columnas:")
            for col in columnas:
                contenido.append(f"  - {col[1]} ({col[2]})")
            
            cursor.execute(f"SELECT * FROM {nombre_tabla} LIMIT 5")
            datos = cursor.fetchall()
            if datos:
                contenido.append("Datos de ejemplo:")
                for fila in datos:
                    contenido.append(f"  {fila}")
            
            contenido.append("")
        
        conn.close()
        return '\n'.join(contenido)
    except Exception as e:
        return f"Error procesando SQLITE: {e}"

def procesar_curso():
    """Procesa todo el curso y extrae la información"""
    
    # Ruta del curso
    curso_path = r"C:\Users\alejo\Escritorio\_Data Science - Machine Learning Flex-20250802T150916Z-1-001\_Data Science - Machine Learning Flex"
    
    if not os.path.exists(curso_path):
        print(f"Error: No se encuentra el directorio del curso en {curso_path}")
        return
    
    # Obtener el directorio donde está este script
    script_dir = os.path.dirname(os.path.abspath(__file__))
    
    print("Procesando curso de Data Science...")
    
    # Diccionario para almacenar todo el contenido
    contenido_completo = {}
    
    # Contadores
    total_archivos = 0
    archivos_procesados = 0
    
    # Procesar archivos recursivamente
    for root, dirs, files in os.walk(curso_path):
        for file in files:
            archivo_completo = os.path.join(root, file)
            extension = os.path.splitext(file)[1].lower()
            
            # Ruta relativa para el diccionario
            ruta_relativa = os.path.relpath(archivo_completo, curso_path)
            
            contenido = ""
            
            # Procesar según la extensión
            if extension == '.docx':
                contenido = extraer_texto_docx(archivo_completo)
                archivos_procesados += 1
            elif extension == '.pdf':
                contenido = extraer_texto_pdf(archivo_completo)
                archivos_procesados += 1
            elif extension == '.ipynb':
                contenido = extraer_texto_ipynb(archivo_completo)
                archivos_procesados += 1
            elif extension == '.pptx':
                contenido = extraer_texto_pptx(archivo_completo)
                archivos_procesados += 1
            elif extension == '.sqlite':
                contenido = extraer_texto_sqlite(archivo_completo)
                archivos_procesados += 1
            
            if contenido:
                contenido_completo[ruta_relativa] = {
                    'tipo': extension[1:],  # Sin el punto
                    'contenido': contenido,
                    'ruta_completa': archivo_completo
                }
            
            total_archivos += 1
            print(f"Procesado: {ruta_relativa}")
    
    print(f"\nProcesamiento completado:")
    print(f"- Archivos totales encontrados: {total_archivos}")
    print(f"- Archivos procesados exitosamente: {archivos_procesados}")
    
    # Guardar en JSON en el directorio del script
    json_path = os.path.join(script_dir, 'contenido_curso_completo.json')
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(contenido_completo, f, ensure_ascii=False, indent=2)
    
    # Crear archivo de texto plano para GPT en el directorio del script
    gpt_path = os.path.join(script_dir, 'contenido_para_gpt.txt')
    with open(gpt_path, 'w', encoding='utf-8') as f:
        f.write("# CONTENIDO COMPLETO DEL CURSO DE DATA SCIENCE\n\n")
        
        for ruta, info in contenido_completo.items():
            f.write(f"## ARCHIVO: {ruta}\n")
            f.write(f"Tipo: {info['tipo']}\n")
            f.write(f"Contenido:\n{info['contenido']}\n")
            f.write("\n" + "="*80 + "\n\n")
    
    # Crear resumen en el directorio del script
    resumen_path = os.path.join(script_dir, 'resumen_curso.txt')
    with open(resumen_path, 'w', encoding='utf-8') as f:
        f.write("# RESUMEN DEL CURSO DE DATA SCIENCE\n\n")
        
        # Contar por tipo
        tipos = {}
        for info in contenido_completo.values():
            tipo = info['tipo']
            tipos[tipo] = tipos.get(tipo, 0) + 1
        
        f.write("## ESTADÍSTICAS:\n")
        for tipo, cantidad in tipos.items():
            f.write(f"- {tipo.upper()}: {cantidad} archivos\n")
        
        f.write(f"\nTotal de archivos procesados: {len(contenido_completo)}\n")
        
        f.write("\n## LISTA DE ARCHIVOS:\n")
        for ruta in sorted(contenido_completo.keys()):
            f.write(f"- {ruta}\n")
    
    print("\nArchivos generados en:", script_dir)
    print("- contenido_curso_completo.json (formato JSON completo)")
    print("- contenido_para_gpt.txt (texto plano para GPT)")
    print("- resumen_curso.txt (resumen y estadísticas)")
    
    print("\n¡Listo! Puedes usar estos archivos para enviar a GPT.")

if __name__ == "__main__":
    procesar_curso() 